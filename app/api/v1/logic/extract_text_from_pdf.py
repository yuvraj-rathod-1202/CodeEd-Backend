from fastapi import HTTPException
from fastapi.responses import JSONResponse
import fitz
from typing import Dict
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract
from io import BytesIO
import speech_recognition as sr
from pydub import AudioSegment
import ffmpeg

def text_cleaning(extracted_text: str) -> Dict[str, str]:
    cleaned_text = extracted_text.replace("\n", " ").replace("\r", " ").strip()
    print("length of text", len(cleaned_text))
    if(len(cleaned_text) == 0):
        raise HTTPException(status_code=500, detail=f"no text found in file")
    if(len(cleaned_text) > 100000):
        raise HTTPException(status_code=500, detail=f"Please upload small file")
    return {"text": cleaned_text}

async def extract_text_from_pdf_logic(pdf_file) -> Dict[str, str]:
    if not pdf_file.filename.endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Only PDF files are allowed.")
    
    file_content = await pdf_file.read()

    try:
        doc = fitz.open(stream=file_content, filetype="pdf")
        extracted_text = ""
        for page in doc:
            extracted_text += page.get_text()
        doc.close()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")
    return text_cleaning(extracted_text)


async def extract_text_from_docx_logic(docx_file) -> Dict[str, str]:
    if not docx_file.filename.endswith(".docx"):
        raise HTTPException(status_code=400, detail="Only DOCX files are allowed.")
    
    file_content = await docx_file.read()

    try:

        document = Document(BytesIO(file_content))
        extracted_text = ""
        for para in document.paragraphs:
            extracted_text += para.text + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing DOCX: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_from_html_logic(html_file) -> Dict[str, str]:
    if not html_file.filename.endswith(".html") and not html_file.filename.endswith(".htm"):
        raise HTTPException(status_code=400, detail="Only HTML files are allowed.")
    
    file_content = await html_file.read()
    try:
        extracted_text = file_content.decode('utf-8')
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing HTML: {str(e)}")

    return text_cleaning(extracted_text)

async def extract_text_from_text_logic(txt_file) -> Dict[str, str]:
    if not txt_file.filename.endswith(".txt"):
        raise HTTPException(status_code=400, detail="Only TXT files are allowed.")
    
    file_content = await txt_file.read()
    try:
        extracted_text = file_content.decode('utf-8')
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing TXT: {str(e)}")

    return text_cleaning(extracted_text)

async def extract_text_from_pptx_logic(pptx_file) -> Dict[str, str]:
    if not pptx_file.filename.endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Only PPTX files are allowed.")
    
    file_content = await pptx_file.read()
    try:
        presentation = Presentation(BytesIO(file_content))
        extracted_text = ""
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    extracted_text += getattr(shape, "text") + "\n"
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PPTX: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_from_md_logic(md_file) -> Dict[str, str]:
    if not md_file.filename.endswith(".md"):
        raise HTTPException(status_code=400, detail="Only MD files are allowed.")
    
    file_content = await md_file.read()
    try:
        extracted_text = file_content.decode('utf-8')
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing MD: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_from_image_logic(image_file) -> Dict[str, str]:
    if not (image_file.filename.endswith(".png") or image_file.filename.endswith(".jpg") or image_file.filename.endswith(".jpeg")):
        raise HTTPException(status_code=400, detail="Only PNG, JPG, and JPEG files are allowed.")
    
    file_content = await image_file.read()
    try:
        image = Image.open(BytesIO(file_content))
        extracted_text = pytesseract.image_to_string(image)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing image: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_from_audio_logic(audio_file) -> Dict[str, str]:
    # Improved file format validation
    supported_formats = [".wav", ".mp3", ".m4a", ".mp4", ".avi", ".flac", ".ogg", ".aac"]
    file_extension = None
    for ext in supported_formats:
        if audio_file.filename.lower().endswith(ext):
            file_extension = ext
            break
    
    if not file_extension:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported audio format. Supported formats: {', '.join(supported_formats)}"
        )
    
    file_content = await audio_file.read()
    
    # Validate file content
    if len(file_content) == 0:
        raise HTTPException(status_code=400, detail="Audio file is empty")
    
    if len(file_content) > 50 * 1024 * 1024:  # 50 MB limit for audio
        raise HTTPException(status_code=400, detail="Audio file too large. Maximum size is 50MB")
    
    try:
        # Step 1: Load audio with better error handling
        try:
            audio_bytes = BytesIO(file_content)
            audio = AudioSegment.from_file(audio_bytes)
        except Exception as load_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to load audio file. Please ensure the file is not corrupted. Error: {str(load_error)}"
            )
        
        # Step 2: Validate audio properties
        if len(audio) == 0:
            raise HTTPException(status_code=400, detail="Audio file contains no audio data")
        
        if len(audio) > 10 * 60 * 1000:  # 10 minutes limit
            raise HTTPException(status_code=400, detail="Audio file too long. Maximum duration is 10 minutes")
        
        # Step 3: Convert audio to proper format for speech recognition
        try:
            # Convert to mono, 16kHz WAV for better speech recognition
            audio = audio.set_frame_rate(16000).set_channels(1)
            
            # Export to WAV format in memory
            temp_wav = BytesIO()
            audio.export(temp_wav, format="wav")
            temp_wav.seek(0)
        except Exception as convert_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to convert audio format. Error: {str(convert_error)}"
            )
        
        # Step 4: Speech recognition with better error handling
        try:
            recognizer = sr.Recognizer()
            
            # Adjust recognizer settings for better accuracy
            recognizer.energy_threshold = 300
            recognizer.dynamic_energy_threshold = True
            recognizer.dynamic_energy_adjustment_damping = 0.15
            recognizer.dynamic_energy_ratio = 1.5
            recognizer.pause_threshold = 0.8
            recognizer.operation_timeout = None
            recognizer.phrase_threshold = 0.3
            recognizer.non_speaking_duration = 0.5
            
            with sr.AudioFile(temp_wav) as source:
                # Adjust for ambient noise
                recognizer.adjust_for_ambient_noise(source, duration=1)
                audio_data = recognizer.record(source)
                
                # Try Google Speech Recognition
                try:
                    # Use getattr to handle dynamic method access
                    recognize_method = getattr(recognizer, 'recognize_google', None)
                    if recognize_method:
                        extracted_text = recognize_method(audio_data, language='en-US')
                    else:
                        raise AttributeError("Google Speech Recognition not available")
                except sr.UnknownValueError:
                    raise HTTPException(
                        status_code=500, 
                        detail="Could not understand audio. Please ensure the audio contains clear speech"
                    )
                except (sr.RequestError, AttributeError) as e:
                    # Fallback to basic recognition without language parameter
                    try:
                        recognize_method = getattr(recognizer, 'recognize_google', None)
                        if recognize_method:
                            extracted_text = recognize_method(audio_data)
                        else:
                            raise HTTPException(
                                status_code=500, 
                                detail="Speech recognition service not available. Please ensure SpeechRecognition library is properly installed"
                            )
                    except (sr.UnknownValueError, sr.RequestError, AttributeError):
                        raise HTTPException(
                            status_code=500, 
                            detail=f"Speech recognition service unavailable. Error: {str(e)}"
                        )
                        
        except HTTPException:
            raise  # Re-raise HTTP exceptions
        except Exception as recognition_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Speech recognition failed. Error: {str(recognition_error)}"
            )
        
        # Validate extracted text
        if not extracted_text or len(extracted_text.strip()) == 0:
            raise HTTPException(
                status_code=500, 
                detail="No speech detected in audio file. Please ensure the audio contains clear speech"
            )
            
    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing audio: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_from_video_logic(video_file) -> Dict[str, str]:
    # Improved file format validation
    supported_formats = [".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv"]
    file_extension = None
    for ext in supported_formats:
        if video_file.filename.lower().endswith(ext):
            file_extension = ext
            break
    
    if not file_extension:
        raise HTTPException(
            status_code=400, 
            detail=f"Unsupported video format. Supported formats: {', '.join(supported_formats)}"
        )
    
    file_content = await video_file.read()
    
    # Validate file content
    if len(file_content) == 0:
        raise HTTPException(status_code=400, detail="Video file is empty")
    
    if len(file_content) > 100 * 1024 * 1024:  # 100 MB limit for video
        raise HTTPException(status_code=400, detail="Video file too large. Maximum size is 100MB")
    
    try:
        # Step 1: Extract audio from video using ffmpeg
        try:
            out, err = (
                ffmpeg
                .input('pipe:0', format='mp4')
                .output('pipe:1', format='wav', acodec='pcm_s16le', ac=1, ar='16000')
                .run(input=file_content, capture_stdout=True, capture_stderr=True)
            )
            
            if err:
                print(f"FFmpeg stderr: {err.decode()}")
            
            if not out:
                raise HTTPException(
                    status_code=500, 
                    detail="Failed to extract audio from video. Please ensure the video contains audio"
                )
                
        except Exception as ffmpeg_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to extract audio from video. Error: {str(ffmpeg_error)}"
            )

        # Step 2: Process extracted audio
        try:
            temp_audio = BytesIO(out)
            
            # Validate extracted audio
            if temp_audio.getvalue().__len__() == 0:
                raise HTTPException(
                    status_code=500, 
                    detail="No audio found in video file"
                )
                
        except Exception as audio_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Failed to process extracted audio. Error: {str(audio_error)}"
            )

        # Step 3: Speech recognition
        try:
            recognizer = sr.Recognizer()
            
            # Adjust recognizer settings
            recognizer.energy_threshold = 300
            recognizer.dynamic_energy_threshold = True
            recognizer.pause_threshold = 0.8
            
            with sr.AudioFile(temp_audio) as source:
                # Adjust for ambient noise
                recognizer.adjust_for_ambient_noise(source, duration=1)
                audio_data = recognizer.record(source)
                
                # Try Google Speech Recognition
                try:
                    recognize_method = getattr(recognizer, 'recognize_google', None)
                    if recognize_method:
                        extracted_text = recognize_method(audio_data, language='en-US')
                    else:
                        raise AttributeError("Google Speech Recognition not available")
                except sr.UnknownValueError:
                    raise HTTPException(
                        status_code=500, 
                        detail="Could not understand audio in video. Please ensure the video contains clear speech"
                    )
                except (sr.RequestError, AttributeError) as e:
                    # Fallback to basic recognition
                    try:
                        recognize_method = getattr(recognizer, 'recognize_google', None)
                        if recognize_method:
                            extracted_text = recognize_method(audio_data)
                        else:
                            raise HTTPException(
                                status_code=500, 
                                detail="Speech recognition service not available"
                            )
                    except (sr.UnknownValueError, sr.RequestError, AttributeError):
                        raise HTTPException(
                            status_code=500, 
                            detail=f"Speech recognition service unavailable. Error: {str(e)}"
                        )
                        
        except HTTPException:
            raise  # Re-raise HTTP exceptions
        except Exception as recognition_error:
            raise HTTPException(
                status_code=500, 
                detail=f"Speech recognition failed. Error: {str(recognition_error)}"
            )
        
        # Validate extracted text
        if not extracted_text or len(extracted_text.strip()) == 0:
            raise HTTPException(
                status_code=500, 
                detail="No speech detected in video. Please ensure the video contains clear speech"
            )
            
    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing video: {str(e)}")
    
    return text_cleaning(extracted_text)

async def extract_text_logic(file) -> Dict[str, str]:
    # Validate file size
    if file.size > 10 * 1024 * 1024:  # 10 MB
        return {"text": "File size exceeds the 10MB limit. Please upload a smaller file."}
    
    # Get file extension
    filename_lower = file.filename.lower()
    
    try:
        if filename_lower.endswith(".pdf"):
            textObj = await extract_text_from_pdf_logic(file)
            return textObj

        if filename_lower.endswith(".txt"):
            textObj = await extract_text_from_text_logic(file)
            return textObj

        if filename_lower.endswith(".docx"):
            textObj = await extract_text_from_docx_logic(file)
            return textObj

        if filename_lower.endswith(".html"):
            textObj = await extract_text_from_html_logic(file)
            return textObj
        
        if filename_lower.endswith(".pptx"):
            textObj = await extract_text_from_pptx_logic(file)
            return textObj
        
        if filename_lower.endswith(".md"):
            textObj = await extract_text_from_md_logic(file)
            return textObj
        
        if filename_lower.endswith((".png", ".jpg", ".jpeg")):
            textObj = await extract_text_from_image_logic(file)
            return textObj
        
        if filename_lower.endswith((".wav", ".mp3", ".m4a", ".flac", ".ogg", ".aac")):
            textObj = await extract_text_from_audio_logic(file)
            return textObj
        
        if filename_lower.endswith((".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv")):
            textObj = await extract_text_from_video_logic(file)
            return textObj

        return {"text": "No valid file format found. Please upload a PDF, DOCX, HTML, TXT, PPTX, IMAGE, Audio (WAV/MP3/M4A/FLAC/OGG/AAC), Video (MP4/AVI/MOV/MKV/WEBM/FLV) or MD file."}
        
    except HTTPException:
        raise  # Re-raise HTTP exceptions
    except Exception as e:
        print(f"Unexpected error in extract_text_logic: {str(e)}")
        raise HTTPException(
            status_code=500, 
            detail=f"An unexpected error occurred while processing the file: {str(e)}"
        )